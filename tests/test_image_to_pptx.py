"""
テスト: 画像 → 要素分析 → 個別SVG → PPTX出力
"""

import os
import sys
import json
import re
from pathlib import Path
from io import BytesIO

sys.path.insert(0, str(Path(__file__).parent.parent))

from dotenv import load_dotenv
from google import genai
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

load_dotenv(dotenv_path=Path(__file__).parent.parent / '.env.local')

OUTPUT_BASE_DIR = Path(__file__).parent / "output"
MODEL = "gemini-3-pro-preview"

# フォントパス
FONT_PATH = Path(__file__).parent.parent / "agents" / "fonts" / "NotoSansCJKjp-Regular.otf"


def generate_test_id() -> str:
    """テストIDを生成"""
    import random
    import string
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    suffix = ''.join(random.choices(string.ascii_uppercase + string.digits, k=4))
    return f"{timestamp}_{suffix}"

# スライドサイズ（16:9）
SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080

PROMPT = """この画像を分析し、編集可能なPowerPointスライドを作成するために、全ての要素を識別してください。

## タスク
1. 画像内の全ての要素を識別する
2. テキスト要素のみSVGを生成する
3. 非テキスト要素（背景、イラスト、写真など）はbbox情報のみ出力

## 出力形式（JSON）
```json
{
  "elements": [
    {
      "id": "background",
      "type": "background",
      "label": "背景の説明",
      "bbox": {"x": 0, "y": 0, "width": 1376, "height": 768}
    },
    {
      "id": "text_1",
      "type": "text",
      "label": "テキストの説明",
      "content": "テキスト内容",
      "bbox": {"x": 0, "y": 0, "width": 100, "height": 50},
      "svg": "<svg>テキストのSVG</svg>"
    },
    {
      "id": "illustration_1",
      "type": "illustration",
      "label": "イラストの説明",
      "bbox": {"x": 0, "y": 0, "width": 100, "height": 100}
    }
  ]
}
```

## ルール
- type="text" の要素のみSVGを生成する
- type="background", "illustration", "image" などはbbox情報のみ（SVGは生成しない）
- bboxは元画像のピクセル座標で指定（この画像は1376x768）
- テキストSVGは元の見た目を忠実に再現（グラデーション、縁取り、影など）
- テキストのfont-familyは "Noto Sans CJK JP" を使用

JSONのみを出力してください。"""


def analyze_image(client, image: Image.Image) -> dict:
    """画像を分析して要素リストを取得"""
    response = client.models.generate_content(
        model=MODEL,
        contents=[PROMPT, image]
    )

    # レスポンス確認
    text = ""
    if response.candidates and response.candidates[0].content:
        for part in response.candidates[0].content.parts:
            if part.text:
                text += part.text

    if not text:
        raise ValueError("APIからテキストレスポンスがありません")

    # JSONを抽出
    json_match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
    if json_match:
        return json.loads(json_match.group(1))

    json_match = re.search(r'\{[\s\S]*\}', text)
    if json_match:
        return json.loads(json_match.group())

    raise ValueError(f"JSONが見つかりません: {text[:500]}")


def save_svg(svg_content: str, filepath: Path):
    """SVGをファイルに保存"""
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(svg_content)


def create_pptx(elements: list, original_image: Image.Image, output_path: Path, output_dir: Path):
    """
    要素リストからPPTXを作成
    - テキスト要素: SVGからPNG変換
    - 非テキスト要素: 元画像から切り出し

    Args:
        elements: 要素リスト
        original_image: 元画像（PIL Image）
        output_path: 出力先パス
        output_dir: 出力ディレクトリ
    """
    from pptx.enum.text import PP_ALIGN
    import cairosvg

    prs = Presentation()

    # スライドサイズを16:9に設定
    prs.slide_width = Emu(SLIDE_WIDTH * 914400 // 96)
    prs.slide_height = Emu(SLIDE_HEIGHT * 914400 // 96)

    # 空白スライドを追加
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 座標変換用スケール
    img_width, img_height = original_image.size
    scale_x = SLIDE_WIDTH / img_width
    scale_y = SLIDE_HEIGHT / img_height

    # 各要素を追加（背景を先に、他の要素を後に）
    sorted_elements = sorted(elements, key=lambda e: 0 if e.get("type") == "background" else 1)

    for elem in sorted_elements:
        bbox = elem.get("bbox", {})
        elem_type = elem.get("type", "")
        content = elem.get("content", "")
        svg = elem.get("svg", "")
        elem_id = elem.get("id", "unknown")

        png_path = output_dir / f"{elem_id}.png"

        # 座標計算
        bx = int(bbox.get("x", 0))
        by = int(bbox.get("y", 0))
        bw = int(bbox.get("width", 100))
        bh = int(bbox.get("height", 50))

        if elem_type == "text" and svg:
            # テキスト要素 → SVGからPNG変換
            svg_path = output_dir / f"{elem_id}.svg"
            save_svg(svg, svg_path)

            try:
                svg_content = svg.replace('Arial Black', 'Noto Sans CJK JP')
                svg_content = svg_content.replace('Impact', 'Noto Sans CJK JP')
                svg_content = svg_content.replace('sans-serif', 'Noto Sans CJK JP')
                cairosvg.svg2png(bytestring=svg_content.encode('utf-8'), write_to=str(png_path))
            except Exception as e:
                print(f"    SVG→PNG変換失敗: {elem_id} - {e}")
                continue
        else:
            # 非テキスト要素 → 元画像から切り出し
            try:
                cropped = original_image.crop((bx, by, bx + bw, by + bh))
                cropped.save(png_path)
            except Exception as e:
                print(f"    画像切り出し失敗: {elem_id} - {e}")
                continue

        # PPTXへの配置座標
        x = Emu(int(bx * scale_x) * 914400 // 96)
        y = Emu(int(by * scale_y) * 914400 // 96)
        width = Emu(int(bw * scale_x) * 914400 // 96)
        height = Emu(int(bh * scale_y) * 914400 // 96)

        if elem_type == "background":
            # 背景 → 全画面に配置
            slide.shapes.add_picture(
                str(png_path),
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )
            print(f"    背景追加: {elem_id}")

        elif elem_type == "text" and content:
            # テキスト要素 → 画像のみ（テキストボックスは追加しない）
            # 編集が必要な場合はユーザーがPPTX上で画像を削除してテキストを入力
            slide.shapes.add_picture(str(png_path), x, y, width, height)

            print(f"    テキスト追加: {elem_id} - {content[:20]}")

        else:
            # その他の要素 → 画像として配置
            slide.shapes.add_picture(str(png_path), x, y, width, height)
            print(f"    要素追加: {elem_id} ({elem_type})")

    # 保存
    prs.save(output_path)
    print(f"\nPPTX保存: {output_path}")


def main():
    """メイン処理"""
    # テストID生成
    test_id = generate_test_id()
    output_dir = OUTPUT_BASE_DIR / test_id
    output_dir.mkdir(parents=True, exist_ok=True)

    print("=== 画像 → SVG → PPTX テスト ===")
    print(f"テストID: {test_id}\n")

    # 入力画像
    image_path = Path(__file__).parent.parent / "output_parse_image" / "IZ99-8459.png"
    print(f"入力画像: {image_path}")

    image = Image.open(image_path)
    print(f"サイズ: {image.size}\n")

    # クライアント初期化
    client = genai.Client(api_key=os.environ.get("GOOGLE_API_KEY"))

    # Step 1: 画像分析
    print("[Step 1] 画像分析中...")
    result = analyze_image(client, image)
    elements = result.get("elements", [])
    print(f"  検出された要素: {len(elements)}個")

    for elem in elements:
        desc = elem.get('content') or elem.get('label') or ''
        print(f"    - {elem['id']}: {elem.get('type')} - {desc[:30]}")

    # 分析結果を保存
    with open(output_dir / "analysis.json", "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"  分析結果保存: {output_dir / 'analysis.json'}")

    # Step 2: 各SVGを保存
    print("\n[Step 2] SVGファイル保存中...")
    for elem in elements:
        svg_path = output_dir / f"{elem['id']}.svg"
        save_svg(elem.get("svg", ""), svg_path)
        print(f"  保存: {svg_path}")

    # Step 3: PPTX作成
    print("\n[Step 3] PPTX作成中...")
    pptx_path = output_dir / "result.pptx"
    create_pptx(elements, image, pptx_path, output_dir)

    print("\n=== 完了 ===")
    print(f"出力ディレクトリ: {output_dir}")


if __name__ == "__main__":
    main()
