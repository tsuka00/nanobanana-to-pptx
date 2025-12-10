"""
SVG to PPTX 変換ツール
SVGをPowerPointプレゼンテーション形式に変換
"""

import base64
import io
import re
from pathlib import Path
from typing import Optional

from strands import tool

# キャンバスサイズ（16:9）
CANVAS_WIDTH = 1920
CANVAS_HEIGHT = 1080

# 出力ディレクトリ
OUTPUT_PPTX_DIR = Path(__file__).parent.parent.parent / "output_pptx"


def svg_to_png_bytes(svg: str) -> bytes:
    """
    SVGをPNGバイトデータに変換
    cairosvgを使用してレンダリング
    """
    try:
        import cairosvg
        png_bytes = cairosvg.svg2png(
            bytestring=svg.encode('utf-8'),
            output_width=CANVAS_WIDTH,
            output_height=CANVAS_HEIGHT
        )
        return png_bytes
    except ImportError:
        raise ImportError("cairosvg is required for SVG to PNG conversion. Install with: pip install cairosvg")


@tool
def svg_to_pptx(
    svg: str,
    session_id: str,
    folder: str = "result"
) -> dict:
    """
    SVGをPowerPointプレゼンテーション（.pptx）に変換して保存します。
    SVGは画像としてレンダリングされ、スライドに配置されます。
    （編集不可の画像として出力）

    Args:
        svg: SVG文字列
        session_id: セッションID（ファイル名に使用）
        folder: 保存先フォルダ名（デフォルト: "result"）

    Returns:
        dict: 変換結果
            - success: 成功したかどうか
            - file_path: 保存されたPPTXファイルのパス
            - error: エラーメッセージ（失敗時）
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx is required. Install with: pip install python-pptx"
        }

    try:
        # 出力ディレクトリを作成
        output_dir = OUTPUT_PPTX_DIR / folder
        output_dir.mkdir(parents=True, exist_ok=True)

        output_path = output_dir / f"{session_id}.pptx"

        # SVGをPNGに変換
        png_bytes = svg_to_png_bytes(svg)

        # PowerPointプレゼンテーションを作成
        prs = Presentation()

        # スライドサイズを16:9に設定（EMU単位）
        # 1インチ = 914400 EMU
        prs.slide_width = Emu(CANVAS_WIDTH * 914400 // 96)  # 96dpi想定
        prs.slide_height = Emu(CANVAS_HEIGHT * 914400 // 96)

        # 空白レイアウトでスライドを追加
        blank_layout = prs.slide_layouts[6]  # 空白レイアウト
        slide = prs.slides.add_slide(blank_layout)

        # PNGをスライドに追加（フル画面）
        image_stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(
            image_stream,
            Emu(0),  # left
            Emu(0),  # top
            prs.slide_width,  # width
            prs.slide_height  # height
        )

        # PPTXを保存
        prs.save(str(output_path))

        return {
            "success": True,
            "file_path": str(output_path)
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }


def _px_to_emu(px: float) -> int:
    """ピクセルをEMU（English Metric Units）に変換（96dpi想定）"""
    return int(px * 914400 / 96)


def _hex_to_rgb(hex_color: str) -> tuple:
    """HEXカラーをRGBタプルに変換"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


@tool
def svg_to_pptx_editable(
    session_id: str,
    folder: str = "result",
    background_base64: Optional[str] = None,
    illustration_base64: Optional[str] = None,
    title_config: Optional[dict] = None,
    subtitle_config: Optional[dict] = None
) -> dict:
    """
    個別の要素からPowerPointプレゼンテーション（.pptx）を生成します。
    テキストはPowerPointのネイティブテキストボックスとして配置され、編集可能です。

    Args:
        session_id: セッションID（ファイル名に使用）
        folder: 保存先フォルダ名（デフォルト: "result"）
        background_base64: 背景画像のBase64データ
        illustration_base64: イラスト画像のBase64データ
        title_config: タイトル設定
            - text: テキスト
            - x, y: 座標
            - fontSize: フォントサイズ
            - fontFamily: フォント名
            - fontWeight: フォントの太さ
            - color: 色（HEX）
        subtitle_config: サブタイトル設定（title_configと同じ形式）

    Returns:
        dict: 変換結果
            - success: 成功したかどうか
            - file_path: 保存されたPPTXファイルのパス
            - error: エラーメッセージ（失敗時）
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx is required. Install with: pip install python-pptx"
        }

    try:
        # 出力ディレクトリを作成
        output_dir = OUTPUT_PPTX_DIR / folder
        output_dir.mkdir(parents=True, exist_ok=True)

        output_path = output_dir / f"{session_id}.pptx"

        # PowerPointプレゼンテーションを作成
        prs = Presentation()

        # スライドサイズを16:9に設定
        prs.slide_width = Emu(_px_to_emu(CANVAS_WIDTH))
        prs.slide_height = Emu(_px_to_emu(CANVAS_HEIGHT))

        # 空白レイアウトでスライドを追加
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 1. 背景画像を追加
        if background_base64:
            bg_bytes = base64.b64decode(background_base64)
            bg_stream = io.BytesIO(bg_bytes)
            slide.shapes.add_picture(
                bg_stream,
                Emu(0),
                Emu(0),
                prs.slide_width,
                prs.slide_height
            )

        # 2. イラストを追加（あれば）
        if illustration_base64:
            illust_bytes = base64.b64decode(illustration_base64)
            illust_stream = io.BytesIO(illust_bytes)
            slide.shapes.add_picture(
                illust_stream,
                Emu(0),
                Emu(0),
                prs.slide_width,
                prs.slide_height
            )

        # 3. タイトルをテキストボックスとして追加
        if title_config and title_config.get("text"):
            _add_text_box(
                slide,
                title_config,
                prs.slide_width,
                prs.slide_height,
                Pt,
                Emu,
                RGBColor,
                PP_ALIGN,
                MSO_ANCHOR
            )

        # 4. サブタイトルをテキストボックスとして追加
        if subtitle_config and subtitle_config.get("text"):
            _add_text_box(
                slide,
                subtitle_config,
                prs.slide_width,
                prs.slide_height,
                Pt,
                Emu,
                RGBColor,
                PP_ALIGN,
                MSO_ANCHOR
            )

        # PPTXを保存
        prs.save(str(output_path))

        return {
            "success": True,
            "file_path": str(output_path)
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }


def _add_text_box(slide, config, slide_width, slide_height, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR):
    """テキストボックスを追加"""
    text = config.get("text", "")
    x = config.get("x", CANVAS_WIDTH // 2)
    y = config.get("y", CANVAS_HEIGHT // 2)
    font_size = config.get("fontSize", 48)
    font_family = config.get("fontFamily", "Arial")
    font_weight = config.get("fontWeight", "normal")
    color = config.get("color", "#FFFFFF")

    # テキストボックスのサイズを計算（幅はスライド幅の80%、中央配置）
    box_width = int(CANVAS_WIDTH * 0.9)
    box_height = int(font_size * 2.5 * (text.count('\n') + 1))

    # 中央揃えを考慮した位置計算
    left = _px_to_emu((CANVAS_WIDTH - box_width) // 2)
    top = _px_to_emu(y - font_size)

    # テキストボックスを追加
    textbox = slide.shapes.add_textbox(
        Emu(left),
        Emu(top),
        Emu(_px_to_emu(box_width)),
        Emu(_px_to_emu(box_height))
    )

    tf = textbox.text_frame
    tf.word_wrap = True

    # 複数行対応
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.alignment = PP_ALIGN.CENTER

        # フォント設定
        run = p.runs[0] if p.runs else p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = font_family
        run.font.bold = font_weight == "bold"

        # 色設定
        try:
            r, g, b = _hex_to_rgb(color)
            run.font.color.rgb = RGBColor(r, g, b)
        except Exception:
            run.font.color.rgb = RGBColor(255, 255, 255)

    # 垂直方向の中央揃え
    tf.anchor = MSO_ANCHOR.MIDDLE
