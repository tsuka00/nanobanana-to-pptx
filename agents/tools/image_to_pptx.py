"""
PPTXを生成するツール
- 背景/イラスト: 生成された画像を配置
- テキスト要素: 編集可能なテキストボックスとして配置
- 図形要素: 四角形、円、線などを配置
"""

import os
import base64
from pathlib import Path
from io import BytesIO
from typing import Optional, List
from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# フォントパス
FONT_PATH = Path(__file__).parent.parent / "fonts" / "NotoSansCJKjp-Regular.otf"

# スライドサイズ（16:9）
SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080

# 出力ベースディレクトリ
AGENT_OUTPUT_DIR = Path(__file__).parent.parent.parent / "agent_output"


def image_to_pptx(
    elements: List[dict],
    session_id: str,
    original_image_base64: Optional[str] = None,
    output_dir: Optional[Path] = None
) -> dict:
    """
    要素リストからPPTXを生成

    Args:
        elements: 要素リスト
            - type="background": 背景画像（image_base64またはfile_pathで指定）
            - type="image": 画像（image_base64またはfile_pathで指定）
            - type="text": テキストボックス（content, style, bboxで指定）
            - type="shape": 図形（shape, bbox, styleで指定）
        session_id: セッションID
        original_image_base64: 非推奨（後方互換性のため残存）
        output_dir: 出力ディレクトリ（省略時はagent_output/{session_id}）

    Returns:
        dict: {
            "success": bool,
            "file_path": str,
            "element_files": list,
            "error": str  # エラー時のみ
        }
    """
    try:
        # 出力ディレクトリ
        out_dir = output_dir or (AGENT_OUTPUT_DIR / session_id)
        out_dir.mkdir(parents=True, exist_ok=True)

        # PPTX作成
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_WIDTH * 914400 // 96)
        prs.slide_height = Emu(SLIDE_HEIGHT * 914400 // 96)

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        element_files = []

        # レイヤー順: 背景 → 画像 → テキスト
        def sort_key(e):
            t = e.get("type", "")
            order = {"background": 0, "image": 1, "text": 2}
            return order.get(t, 1)

        sorted_elements = sorted(elements, key=sort_key)

        for elem in sorted_elements:
            elem_type = elem.get("type", "")
            elem_id = elem.get("id", "unknown")

            if elem_type == "text":
                # テキスト要素 → 編集可能なテキストボックス
                _add_textbox(slide, elem, prs)

            elif elem_type == "background":
                # 背景画像 → 全画面配置
                png_path = _get_image_path(elem, out_dir, elem_id)
                if png_path:
                    element_files.append(str(png_path))
                    slide.shapes.add_picture(
                        str(png_path),
                        Emu(0), Emu(0),
                        prs.slide_width, prs.slide_height
                    )

            elif elem_type == "image":
                # 画像要素 → bbox位置に配置
                png_path = _get_image_path(elem, out_dir, elem_id)
                if png_path:
                    element_files.append(str(png_path))
                    bbox = elem.get("bbox", {})
                    x, y, width, height = _bbox_to_emu(bbox)
                    slide.shapes.add_picture(str(png_path), x, y, width, height)

        # 保存
        pptx_path = out_dir / f"{session_id}.pptx"
        prs.save(pptx_path)

        return {
            "success": True,
            "file_path": str(pptx_path),
            "element_files": element_files
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }


def _get_image_path(elem: dict, out_dir: Path, elem_id: str) -> Optional[Path]:
    """要素から画像パスを取得（file_pathまたはimage_base64から）"""
    # 既にファイルパスがある場合
    if elem.get("file_path"):
        return Path(elem["file_path"])

    # Base64から保存
    if elem.get("image_base64"):
        png_path = out_dir / f"{elem_id}.png"
        image_data = base64.b64decode(elem["image_base64"])
        image = Image.open(BytesIO(image_data))
        image.save(png_path)
        return png_path

    return None


def _bbox_to_emu(bbox: dict) -> tuple:
    """bboxをEMU単位の座標に変換"""
    x = int(bbox.get("x", 0))
    y = int(bbox.get("y", 0))
    w = int(bbox.get("width", 100))
    h = int(bbox.get("height", 100))

    # ピクセル → EMU変換（96dpi基準）
    return (
        Emu(x * 914400 // 96),
        Emu(y * 914400 // 96),
        Emu(w * 914400 // 96),
        Emu(h * 914400 // 96)
    )


def _add_textbox(slide, elem: dict, prs) -> None:
    """テキストボックスを追加"""
    content = elem.get("content", "")
    style = elem.get("style", {})
    bbox = elem.get("bbox", {})

    x, y, width, height = _bbox_to_emu(bbox)

    textbox = slide.shapes.add_textbox(x, y, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content

    # フォントサイズ
    font_size = style.get("fontSize", 24)
    run.font.size = Pt(font_size)

    # フォント名
    run.font.name = "Noto Sans CJK JP"

    # フォントウェイト
    if style.get("fontWeight") == "bold":
        run.font.bold = True

    # フォントスタイル
    if style.get("fontStyle") == "italic":
        run.font.italic = True

    # 色
    color = style.get("color", "#000000")
    if color.startswith("#"):
        color = color[1:]
    try:
        r = int(color[0:2], 16)
        g = int(color[2:4], 16)
        b = int(color[4:6], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except (ValueError, IndexError):
        pass

    # 配置
    align = style.get("align", "left")
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
